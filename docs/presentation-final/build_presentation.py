"""Build the DSPRO1 Final Presentation (pitch format, 20 min, ~22 slides).

Run from the project root:
    pip install python-pptx
    python docs/presentation-final/build_presentation.py

The script overwrites
    docs/presentation-final/DISPRO1_PresentationFinal_Team8_PredictingApartmentRentalPrices.pptx

The deck mirrors the structure of the Final Report and reuses the HSLU style
of the Mid-Term presentation (HSLU red accent, dark sandwich slides for title
and section dividers, light slides for content).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PATH = PROJECT_ROOT / "docs" / "presentation-final" / (
    "DISPRO1_PresentationFinal_Team8_PredictingApartmentRentalPrices.pptx"
)
FIG_DIR = PROJECT_ROOT / "docs" / "final-report" / "fig"

# Slide size: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# HSLU palette
HSLU_RED       = RGBColor(0xE3, 0x06, 0x13)
DARK           = RGBColor(0x14, 0x14, 0x18)
SOFT_DARK      = RGBColor(0x2A, 0x2A, 0x2F)
TEXT_DARK      = RGBColor(0x22, 0x22, 0x22)
TEXT_GRAY      = RGBColor(0x55, 0x55, 0x55)
TEXT_LIGHT     = RGBColor(0xF2, 0xF2, 0xF2)
ACCENT_BLUE    = RGBColor(0x1F, 0x4E, 0x79)
GREEN_OK       = RGBColor(0x3E, 0xA0, 0x6F)
AMBER          = RGBColor(0xE8, 0x9C, 0x1A)
CARD_BG        = RGBColor(0xF7, 0xF7, 0xF8)
CARD_BORDER    = RGBColor(0xE3, 0xE3, 0xE6)

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_blank_slide(prs: Presentation, dark: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK if dark else RGBColor(0xFF, 0xFF, 0xFF)
    return slide


def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, italic=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, *,
                size=16, color=TEXT_DARK, bullet_color=HSLU_RED,
                line_spacing=1.15, font=FONT_BODY):
    """Render a bullet list using a coloured square dot prefix."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        dot = p.add_run()
        dot.text = "■  "
        dot.font.name = font
        dot.font.size = Pt(size)
        dot.font.color.rgb = bullet_color
        body = p.add_run()
        body.text = item
        body.font.name = font
        body.font.size = Pt(size)
        body.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill=CARD_BG, line=None,
             rounded=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_image_or_placeholder(slide, fig_filename, left, top, width, height,
                              caption=None):
    """Insert PNG from fig/ if present, otherwise show a labelled placeholder."""
    fig_path = FIG_DIR / fig_filename
    if fig_path.exists():
        slide.shapes.add_picture(str(fig_path), left, top,
                                  width=width, height=height)
    else:
        rect = add_rect(slide, left, top, width, height,
                        fill=CARD_BG, line=CARD_BORDER)
        ph = rect.text_frame
        ph.margin_left = Emu(0); ph.margin_right = Emu(0)
        ph.margin_top = Emu(0); ph.margin_bottom = Emu(0)
        ph.vertical_anchor = MSO_ANCHOR.MIDDLE
        ph.word_wrap = True
        p = ph.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[Figure placeholder]\n{fig_filename}"
        r.font.name = FONT_BODY
        r.font.size = Pt(11)
        r.font.italic = True
        r.font.color.rgb = TEXT_GRAY
    if caption:
        cap_top = top + height + Inches(0.08)
        add_text(slide, caption, left, cap_top, width, Inches(0.4),
                 size=10, italic=True, color=TEXT_GRAY,
                 align=PP_ALIGN.CENTER)


def page_chrome(slide, title, subtitle=None, page=None, total=None):
    """Light slide header: title bar + page indicator at top-right."""
    # Left red strip as visual motif
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H,
             fill=HSLU_RED)
    # Title
    add_text(slide, title, Inches(0.55), Inches(0.35),
             Inches(10.5), Inches(0.7),
             size=28, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, Inches(0.55), Inches(0.95),
                 Inches(10.5), Inches(0.4),
                 size=13, color=TEXT_GRAY, italic=True)
    # Page indicator
    if page is not None and total is not None:
        add_text(slide, f"{page} / {total}",
                 Inches(11.9), Inches(0.4),
                 Inches(1.0), Inches(0.35),
                 size=10, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)
    # Footer
    add_text(slide,
             "DSPRO1 FS26  |  Team 8  |  Predicting Apartment Rental Prices in Switzerland",
             Inches(0.55), Inches(7.05), Inches(12.0), Inches(0.3),
             size=9, color=TEXT_GRAY)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_title(prs):
    s = add_blank_slide(prs, dark=True)
    add_rect(s, Inches(0), Inches(6.7), SLIDE_W, Inches(0.05),
             fill=HSLU_RED)
    add_text(s, "PREDICTING APARTMENT",
             Inches(0.7), Inches(2.2), Inches(12), Inches(1.0),
             size=46, bold=True, color=TEXT_LIGHT, font=FONT_HEAD)
    add_text(s, "RENTAL PRICES IN SWITZERLAND",
             Inches(0.7), Inches(3.0), Inches(12), Inches(1.0),
             size=46, bold=True, color=TEXT_LIGHT, font=FONT_HEAD)
    add_text(s, "A supervised machine-learning approach for the Swiss cold-rent market",
             Inches(0.7), Inches(4.05), Inches(12), Inches(0.5),
             size=18, italic=True, color=RGBColor(0xBB, 0xBB, 0xBE))
    add_rect(s, Inches(0.7), Inches(4.85), Inches(0.5), Inches(0.04),
             fill=HSLU_RED)
    add_text(s, "DSPRO1 Final Presentation  |  Spring 2026 (FS26)",
             Inches(0.7), Inches(5.05), Inches(12), Inches(0.4),
             size=14, color=RGBColor(0xCC, 0xCC, 0xCF))
    add_text(s, "Team 8  |  Elias Martinelli  |  Timo Schlumpf",
             Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
             size=14, color=RGBColor(0xCC, 0xCC, 0xCF))
    add_text(s, "Supervisor: Elena Nazarenko",
             Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
             size=12, color=TEXT_GRAY)


def slide_agenda(prs, total):
    s = add_blank_slide(prs)
    page_chrome(s, "Agenda", "20 minutes  |  six chapters  |  pitch + technical deep-dive",
                page=2, total=total)
    items = [
        ("1", "The problem", "Why is a fair rent so hard to judge?"),
        ("2", "Data and scraping pipeline", "From 25'002 URLs to a clean modelling set"),
        ("3", "Modelling approach", "Six models on the same split"),
        ("4", "Results and error analysis", "What works, what doesn't, and why"),
        ("5", "Demo", "Streamlit app with address lookup"),
        ("6", "Conclusions and next steps", "Where we landed and what's next"),
    ]
    col_w = Inches(6.0)
    row_h = Inches(0.85)
    for i, (num, title, sub) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(1.6 + row * 1.6)
        # Number badge
        badge = add_rect(s, x, y, Inches(0.85), Inches(0.85),
                         fill=HSLU_RED, rounded=True)
        btf = badge.text_frame
        btf.margin_left = Emu(0); btf.margin_right = Emu(0)
        btf.margin_top = Emu(0); btf.margin_bottom = Emu(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = num
        br.font.name = FONT_HEAD; br.font.size = Pt(26)
        br.font.bold = True; br.font.color.rgb = TEXT_LIGHT
        add_text(s, title, x + Inches(1.05), y - Inches(0.02),
                 Inches(5.0), Inches(0.45),
                 size=18, bold=True, color=TEXT_DARK, font=FONT_HEAD)
        add_text(s, sub, x + Inches(1.05), y + Inches(0.42),
                 Inches(5.0), Inches(0.4),
                 size=12, italic=True, color=TEXT_GRAY)


def slide_section_divider(prs, number, title, subtitle, total, page):
    s = add_blank_slide(prs, dark=True)
    add_rect(s, Inches(0.7), Inches(2.5), Inches(0.7), Inches(0.05),
             fill=HSLU_RED)
    add_text(s, f"0{number}", Inches(0.7), Inches(1.4),
             Inches(4.0), Inches(1.4),
             size=110, bold=True, color=HSLU_RED, font=FONT_HEAD)
    add_text(s, title.upper(), Inches(0.7), Inches(2.85),
             Inches(12.0), Inches(0.9),
             size=42, bold=True, color=TEXT_LIGHT, font=FONT_HEAD)
    add_text(s, subtitle, Inches(0.7), Inches(3.85),
             Inches(12.0), Inches(0.5),
             size=18, italic=True, color=RGBColor(0xBB, 0xBB, 0xBE))
    add_text(s, f"{page} / {total}",
             Inches(11.9), Inches(7.05), Inches(1.0), Inches(0.3),
             size=10, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)


# ---- Section 1: Problem -----------------------------------------------------
def slide_problem(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "The problem", page=page, total=total)
    # Left: listing card
    add_rect(s, Inches(0.6), Inches(1.55), Inches(4.6), Inches(4.4),
             fill=DARK, rounded=True)
    add_text(s, "You find this listing",
             Inches(0.85), Inches(1.8), Inches(4.2), Inches(0.4),
             size=13, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s, "3 rooms\n80 m²\nZurich",
             Inches(0.85), Inches(2.25), Inches(4.2), Inches(1.6),
             size=20, color=TEXT_LIGHT, font=FONT_HEAD)
    add_rect(s, Inches(0.85), Inches(4.0), Inches(2.5), Inches(0.04),
             fill=HSLU_RED)
    add_text(s, "CHF 2'400 / month",
             Inches(0.85), Inches(4.15), Inches(4.2), Inches(0.6),
             size=28, bold=True, color=HSLU_RED, font=FONT_HEAD)
    add_text(s, "Is this fair?",
             Inches(0.85), Inches(4.85), Inches(4.2), Inches(0.5),
             size=16, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))

    # Right: reasons
    add_text(s, "Why is this hard to judge?",
             Inches(5.6), Inches(1.7), Inches(7.0), Inches(0.5),
             size=20, bold=True, color=TEXT_DARK, font=FONT_HEAD)

    reasons = [
        ("Location dependency",
         "The same apartment can cost twice as much across cantons."),
        ("Many interacting features",
         "Living area, rooms, floor, year built, amenities, public transport …"),
        ("No transparent benchmark",
         "There is no public per-listing reference model on the market."),
    ]
    for i, (head, body) in enumerate(reasons):
        ytop = Inches(2.35 + i * 1.05)
        add_rect(s, Inches(5.6), ytop, Inches(7.2), Inches(0.95),
                 fill=CARD_BG, rounded=True)
        add_rect(s, Inches(5.6), ytop, Inches(0.12), Inches(0.95),
                 fill=HSLU_RED)
        add_text(s, head, Inches(5.85), ytop + Inches(0.08),
                 Inches(6.8), Inches(0.4),
                 size=14, bold=True, color=TEXT_DARK)
        add_text(s, body, Inches(5.85), ytop + Inches(0.45),
                 Inches(6.8), Inches(0.45),
                 size=12, color=TEXT_GRAY)

    # Bottom takeaway
    add_rect(s, Inches(0.6), Inches(6.15), Inches(12.2), Inches(0.6),
             fill=DARK, rounded=True)
    add_text(s,
             "Our goal: build a supervised ML model that predicts fair rent from listing data.",
             Inches(0.85), Inches(6.27), Inches(12.0), Inches(0.4),
             size=14, bold=True, color=TEXT_LIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_relevance(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Why does it matter?",
                "Three stakeholder groups benefit from a transparent estimate",
                page=page, total=total)
    cards = [
        ("Tenants",
         "Spot overpriced listings before signing the contract; budget realistic rent ranges before moving.",
         HSLU_RED),
        ("Landlords",
         "Set a market-based price for vacancies; calibrate offers without paying for a private appraisal.",
         ACCENT_BLUE),
        ("Society",
         "A transparent reference value supports the discussion around rent affordability and regional fairness.",
         GREEN_OK),
    ]
    card_w = Inches(3.95)
    for i, (head, body, color) in enumerate(cards):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(1.9), card_w, Inches(4.2),
                 fill=CARD_BG, rounded=True)
        add_rect(s, x, Inches(1.9), card_w, Inches(0.5),
                 fill=color, rounded=False)
        add_text(s, head, x + Inches(0.25), Inches(2.0),
                 card_w - Inches(0.5), Inches(0.4),
                 size=15, bold=True, color=TEXT_LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, x + Inches(0.3), Inches(2.7),
                 card_w - Inches(0.6), Inches(3.0),
                 size=14, color=TEXT_DARK)


def slide_hypothesis(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Hypothesis", page=page, total=total)
    add_rect(s, Inches(1.2), Inches(2.0), Inches(11.0), Inches(3.6),
             fill=CARD_BG, rounded=True)
    add_rect(s, Inches(1.2), Inches(2.0), Inches(0.14), Inches(3.6),
             fill=HSLU_RED)
    quote = (
        '"Location-related variables combined with structural apartment\n'
        "features contain sufficient signal to predict monthly rental\n"
        "prices in Switzerland with practically useful accuracy, and a\n"
        "Random Forest Regressor will outperform a Linear Regression\n"
        "model in doing so."
    )
    add_text(s, quote, Inches(1.7), Inches(2.4), Inches(10.4), Inches(2.6),
             size=20, italic=True, color=TEXT_DARK, font=FONT_HEAD,
             align=PP_ALIGN.LEFT)
    add_text(s, "Source: project proposal, mid-term presentation slide 6",
             Inches(1.7), Inches(5.05), Inches(10.4), Inches(0.4),
             size=11, color=TEXT_GRAY, italic=True)
    add_text(s,
             "We test this with a baseline, a linear benchmark, and four tree-based models.",
             Inches(1.2), Inches(6.0), Inches(11.0), Inches(0.5),
             size=15, color=TEXT_DARK, align=PP_ALIGN.CENTER)


# ---- Section 2: Data --------------------------------------------------------
def slide_data_source(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Data source", "Rentumo Rental Listings Switzerland v1.0  |  13 April 2026",
                page=page, total=total)

    # Big number panel
    add_rect(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.4),
             fill=DARK, rounded=True)
    add_text(s, "Single-shot scrape", Inches(0.85), Inches(1.95),
             Inches(5.4), Inches(0.4),
             size=14, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s, "25'002", Inches(0.85), Inches(2.4),
             Inches(5.4), Inches(1.0),
             size=58, bold=True, color=TEXT_LIGHT, font=FONT_HEAD)
    add_text(s, "URLs identified on rentumo.ch",
             Inches(0.85), Inches(3.4), Inches(5.4), Inches(0.4),
             size=12, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_rect(s, Inches(0.85), Inches(3.95), Inches(1.5), Inches(0.04),
             fill=HSLU_RED)
    add_text(s, "9'994", Inches(0.85), Inches(4.1),
             Inches(5.4), Inches(0.9),
             size=44, bold=True, color=HSLU_RED, font=FONT_HEAD)
    add_text(s, "listings successfully extracted",
             Inches(0.85), Inches(4.95), Inches(5.4), Inches(0.4),
             size=12, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s,
             "The remaining 15'008 pages are empty placeholder URLs.\nThis is a property of the platform, not a scraping defect.",
             Inches(0.85), Inches(5.4), Inches(5.4), Inches(0.7),
             size=11, italic=True, color=RGBColor(0xAA, 0xAA, 0xAE))

    # Right panel: bullet facts
    add_text(s, "Why Rentumo, not Comparis?",
             Inches(6.8), Inches(1.85), Inches(6.0), Inches(0.4),
             size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "Comparis requires written authorisation for scraping.",
        "Comparis access cap is around 300 records, not enough for ML.",
        "Rentumo was contacted by e-mail before scraping; no reply was received.",
        "Data is used only for non-commercial academic purposes (HSLU DSPRO1).",
    ], Inches(6.8), Inches(2.45), Inches(6.2), Inches(2.8), size=13)

    add_rect(s, Inches(6.8), Inches(5.2), Inches(6.0), Inches(0.9),
             fill=CARD_BG, rounded=True)
    add_text(s, "Single-platform bias",
             Inches(7.0), Inches(5.27), Inches(5.8), Inches(0.35),
             size=12, bold=True, color=TEXT_DARK)
    add_text(s,
             "All listings come from rentumo.ch. Adding a second platform is the most\nobvious next step (see future work).",
             Inches(7.0), Inches(5.55), Inches(5.8), Inches(0.5),
             size=11, color=TEXT_GRAY, italic=True)


def slide_scraper(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Scraper architecture",
                "ScrapeGoat: custom Rust library, two-phase pipeline",
                page=page, total=total)

    # Phase 1
    add_rect(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(2.4),
             fill=CARD_BG, rounded=True)
    add_rect(s, Inches(0.6), Inches(1.8), Inches(0.14), Inches(2.4),
             fill=HSLU_RED)
    add_text(s, "Phase 1: Listing index pages",
             Inches(0.9), Inches(1.95), Inches(5.6), Inches(0.4),
             size=15, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "750 paginated pages /mietobjekte?page={1..750}",
        "CSS selectors on #listings container",
        "Extracts: rooms, area, cold rent, slug",
        "INSERT ... ON CONFLICT (slug) DO NOTHING",
    ], Inches(0.9), Inches(2.45), Inches(5.5), Inches(1.75), size=12)

    # Phase 2
    add_rect(s, Inches(0.6), Inches(4.4), Inches(6.0), Inches(2.4),
             fill=CARD_BG, rounded=True)
    add_rect(s, Inches(0.6), Inches(4.4), Inches(0.14), Inches(2.4),
             fill=ACCENT_BLUE)
    add_text(s, "Phase 2: Detail pages",
             Inches(0.9), Inches(4.55), Inches(5.6), Inches(0.4),
             size=15, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "One call per slug from the database",
        "Adds address, Nebenkosten, Kaution, description",
        "JSON-LD Apartment schema + Datum verfügbar",
        "INSERT ... ON CONFLICT (listing_id) DO UPDATE",
    ], Inches(0.9), Inches(5.05), Inches(5.5), Inches(1.75),
        size=12, bullet_color=ACCENT_BLUE)

    # Right side: orchestration + politeness
    add_text(s, "Orchestration: Mgt",
             Inches(7.0), Inches(1.95), Inches(6.0), Inches(0.4),
             size=15, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "Up to 10 concurrent in-flight HTTP requests",
        "Channel-based feedback: new request only after one completes",
        "Random User-Agent per request from preloaded list",
        "Optional HTTP proxies via plain-text file",
        "Both phases idempotent: safe to resume or replay",
    ], Inches(7.0), Inches(2.45), Inches(5.8), Inches(2.5),
        size=12, bullet_color=GREEN_OK)

    add_rect(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.8),
             fill=DARK, rounded=True)
    add_text(s, "Bottom line", Inches(7.2), Inches(5.15),
             Inches(5.6), Inches(0.4),
             size=13, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s,
             "Bounded server load, full reproducibility,\nresumable on failure. The scraper itself is a\nproject artefact, not a one-off script.",
             Inches(7.2), Inches(5.55), Inches(5.6), Inches(1.2),
             size=13, color=TEXT_LIGHT)


def slide_enrichment(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Data enrichment",
                "From an address to a feature row that the model can use",
                page=page, total=total)
    steps = [
        ("Address",
         "Listing address from rentumo.ch detail page",
         RGBColor(0x55, 0x55, 0x55)),
        ("Geocode + EGID",
         "GeoAdmin SearchServer returns LV95 coordinates and the EGID",
         HSLU_RED),
        ("GWR (BFS)",
         "Building data: year_built, ganzwhg, garea",
         ACCENT_BLUE),
        ("swisstopo / GeoAdmin",
         "elevation, oev_score, solar_class (1-5), population (hectare grid)",
         GREEN_OK),
        ("Modelling row",
         "12 columns ready for the train/eval split",
         AMBER),
    ]
    box_w = Inches(2.3)
    gap = Inches(0.2)
    total_w = box_w * 5 + gap * 4
    start_x = (SLIDE_W - total_w) / 2
    for i, (head, body, color) in enumerate(steps):
        x = start_x + (box_w + gap) * i
        add_rect(s, x, Inches(2.3), box_w, Inches(3.1),
                 fill=CARD_BG, rounded=True)
        add_rect(s, x, Inches(2.3), box_w, Inches(0.5),
                 fill=color, rounded=False)
        add_text(s, head, x + Inches(0.15), Inches(2.39),
                 box_w - Inches(0.3), Inches(0.35),
                 size=13, bold=True, color=TEXT_LIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, x + Inches(0.2), Inches(2.95),
                 box_w - Inches(0.4), Inches(2.4),
                 size=12, color=TEXT_DARK)
        # Arrow (except last)
        if i < len(steps) - 1:
            arrow_x = x + box_w + Inches(0.01)
            add_text(s, "▶", arrow_x, Inches(3.6),
                     gap, Inches(0.5),
                     size=18, color=TEXT_GRAY,
                     align=PP_ALIGN.CENTER)
    add_text(s,
             "All enrichment is point-in-time on 13 April 2026. The pipeline is documented end-to-end in final_records.ipynb.",
             Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.4),
             size=12, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)


def slide_data_funnel(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "From raw URLs to a usable dataset",
                "Five filters reduce 25'002 candidates to about 4'500 modelling rows",
                page=page, total=total)
    rows = [
        ("URLs on rentumo.ch", "25'002", DARK, 12.0),
        ("Successfully extracted", "9'994", SOFT_DARK, 7.2),
        ("After mandatory-field filter (address, price, area, rooms)", "~5'000", HSLU_RED, 5.0),
        ("After GWR enrichment join", "~4'700", ACCENT_BLUE, 4.5),
        ("Final modelling set", "~4'500", GREEN_OK, 4.3),
    ]
    y0 = 1.9
    bar_h = 0.7
    gap = 0.18
    for i, (label, value, color, width_in) in enumerate(rows):
        y = Inches(y0 + i * (bar_h + gap))
        # Bar
        add_rect(s, Inches(0.6), y, Inches(width_in), Inches(bar_h),
                 fill=color, rounded=True)
        # Label inside bar
        add_text(s, label,
                 Inches(0.85), y + Inches(0.15),
                 Inches(width_in - 0.5), Inches(0.4),
                 size=13, bold=True, color=TEXT_LIGHT,
                 anchor=MSO_ANCHOR.TOP)
        # Big number to the right
        add_text(s, value,
                 Inches(0.7 + width_in + 0.15),
                 y + Inches(0.08),
                 Inches(3.0), Inches(0.55),
                 size=24, bold=True, color=color, font=FONT_HEAD)
    add_text(s,
             "Each step is documented in final_records.ipynb (cumulative_filters.csv).",
             Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.4),
             size=12, italic=True, color=TEXT_GRAY,
             align=PP_ALIGN.CENTER)


# ---- Section 3: Methods -----------------------------------------------------
def slide_methods(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Modelling approach",
                "Six models, one shared split, same evaluation protocol",
                page=page, total=total)
    add_image_or_placeholder(s, "geo_clusters.png",
                              Inches(0.6), Inches(1.85),
                              Inches(5.4), Inches(4.6),
                              caption="KMeans geo-clusters fitted on train_df only.")
    add_text(s, "Model pool", Inches(6.4), Inches(1.85),
             Inches(6.5), Inches(0.4),
             size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "DummyRegressor (median): lower bound",
        "Ridge with StandardScaler: linear benchmark",
        "RandomForestRegressor",
        "GradientBoostingRegressor",
        "XGBoost",
        "LightGBM (best single model)",
        "StackingRegressor with Ridge meta-learner",
    ], Inches(6.4), Inches(2.3), Inches(6.4), Inches(3.1), size=13)
    add_rect(s, Inches(6.4), Inches(5.5), Inches(6.5), Inches(1.0),
             fill=DARK, rounded=True)
    add_text(s, "Feature engineering",
             Inches(6.6), Inches(5.6), Inches(6.2), Inches(0.4),
             size=12, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s,
             "building_age, area_per_room, land_area_per_apartment, geo_cluster, knn_price_mean, knn_price_median",
             Inches(6.6), Inches(5.95), Inches(6.2), Inches(0.5),
             size=12, color=TEXT_LIGHT)


def slide_eval_protocol(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Evaluation protocol",
                "Single split with safety nets: cross-validation and a held-out test set",
                page=page, total=total)

    # Split visualization
    add_text(s, "Hold-out test split (60/20/20)",
             Inches(0.6), Inches(1.8), Inches(8.0), Inches(0.4),
             size=14, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_rect(s, Inches(0.6), Inches(2.25), Inches(7.6), Inches(0.7),
             fill=ACCENT_BLUE, rounded=False)
    add_rect(s, Inches(8.2), Inches(2.25), Inches(2.4), Inches(0.7),
             fill=AMBER, rounded=False)
    add_rect(s, Inches(10.6), Inches(2.25), Inches(2.2), Inches(0.7),
             fill=GREEN_OK, rounded=False)
    add_text(s, "Train  60%", Inches(0.6), Inches(2.35),
             Inches(7.6), Inches(0.5),
             size=15, bold=True, color=TEXT_LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Val  20%", Inches(8.2), Inches(2.35),
             Inches(2.4), Inches(0.5),
             size=14, bold=True, color=TEXT_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, "Test  20%", Inches(10.6), Inches(2.35),
             Inches(2.2), Inches(0.5),
             size=14, bold=True, color=TEXT_LIGHT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s,
             "Test is touched only once, at the very end.",
             Inches(0.6), Inches(3.1), Inches(12.0), Inches(0.3),
             size=11, italic=True, color=TEXT_GRAY)

    # Metrics
    metrics = [
        ("MAE", "Average absolute error in CHF"),
        ("RMSE", "Quadratic error, penalises large misses"),
        ("R²", "Variance explained by the model"),
        ("MedAE", "Robust median absolute error"),
        ("CV std", "Stability across folds"),
        ("Bootstrap CI", "Significance test for model differences"),
    ]
    for i, (head, body) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(3.85 + row * 1.35)
        add_rect(s, x, y, Inches(4.0), Inches(1.15),
                 fill=CARD_BG, rounded=True)
        add_rect(s, x, y, Inches(0.14), Inches(1.15),
                 fill=HSLU_RED)
        add_text(s, head, x + Inches(0.25), y + Inches(0.1),
                 Inches(3.6), Inches(0.4),
                 size=14, bold=True, color=TEXT_DARK)
        add_text(s, body, x + Inches(0.25), y + Inches(0.5),
                 Inches(3.6), Inches(0.65),
                 size=11, color=TEXT_GRAY)


# ---- Section 4: Results -----------------------------------------------------
def slide_baseline(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Baseline: mean-prediction strategy",
                "Predicts the same training-set mean for every apartment, regardless of features",
                page=page, total=total)
    boxes = [
        ("MAE", "CHF 612", "Average prediction error", HSLU_RED),
        ("RMSE", "CHF 847", "Penalises large errors more", ACCENT_BLUE),
        ("R²", "0.00", "No variance explained", DARK),
    ]
    for i, (label, value, sub, color) in enumerate(boxes):
        x = Inches(0.6 + i * 4.25)
        add_rect(s, x, Inches(2.1), Inches(4.0), Inches(2.7),
                 fill=CARD_BG, rounded=True)
        add_rect(s, x, Inches(2.1), Inches(4.0), Inches(0.5),
                 fill=color)
        add_text(s, label, x, Inches(2.18),
                 Inches(4.0), Inches(0.4),
                 size=15, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, value, x, Inches(2.85),
                 Inches(4.0), Inches(1.0),
                 size=44, bold=True, color=color, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sub, x, Inches(4.0),
                 Inches(4.0), Inches(0.7),
                 size=12, color=TEXT_GRAY,
                 align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(1.0),
             fill=DARK, rounded=True)
    add_text(s,
             "Any serious model must clearly beat this threshold.\nWe set RMSE < 800 CHF as our minimum success criterion.",
             Inches(0.85), Inches(5.55), Inches(11.8), Inches(0.8),
             size=14, color=TEXT_LIGHT, anchor=MSO_ANCHOR.MIDDLE)


def slide_model_comparison(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Model comparison",
                "All models on the same split, FEATURES_ENGINEERED",
                page=page, total=total)
    add_image_or_placeholder(s, "rmse_train_eval.png",
                              Inches(0.6), Inches(1.85),
                              Inches(7.5), Inches(4.6),
                              caption="Train vs Eval RMSE per model (Notebook Ch. 11).")
    # Table on the right
    add_text(s, "RMSE Eval (CHF)",
             Inches(8.4), Inches(1.85), Inches(4.6), Inches(0.4),
             size=15, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    rows = [
        ("LightGBM", "399", HSLU_RED, True),
        ("XGBoost", "421", SOFT_DARK, False),
        ("GradientBoosting", "425", SOFT_DARK, False),
        ("RandomForest", "425", SOFT_DARK, False),
        ("Ridge (linear)", "560", SOFT_DARK, False),
        ("Dummy (mean)", "847", SOFT_DARK, False),
    ]
    for i, (name, rmse, color, highlight) in enumerate(rows):
        y = Inches(2.4 + i * 0.55)
        bg = HSLU_RED if highlight else CARD_BG
        add_rect(s, Inches(8.4), y, Inches(4.6), Inches(0.48),
                 fill=bg, rounded=False)
        add_text(s, name, Inches(8.55), y + Inches(0.08),
                 Inches(3.0), Inches(0.4),
                 size=13, bold=highlight,
                 color=TEXT_LIGHT if highlight else TEXT_DARK)
        add_text(s, rmse, Inches(11.4), y + Inches(0.08),
                 Inches(1.5), Inches(0.4),
                 size=14, bold=True,
                 color=TEXT_LIGHT if highlight else TEXT_DARK,
                 align=PP_ALIGN.RIGHT)
    add_text(s,
             "After adding KNN-distance features, LightGBM reaches RMSE 393 / R² 0.751.",
             Inches(8.4), Inches(5.85), Inches(4.6), Inches(0.6),
             size=11, italic=True, color=TEXT_GRAY)


def slide_hypothesis_confirmed(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Hypothesis: confirmed",
                "Tree-based models clearly outperform linear regression",
                page=page, total=total)
    # Before / After style
    add_rect(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(4.7),
             fill=CARD_BG, rounded=True)
    add_text(s, "Mid-term baseline",
             Inches(0.85), Inches(2.0), Inches(5.5), Inches(0.4),
             size=13, italic=True, color=TEXT_GRAY)
    add_text(s, "CHF 847",
             Inches(0.85), Inches(2.55), Inches(5.5), Inches(1.2),
             size=66, bold=True, color=DARK, font=FONT_HEAD)
    add_text(s, "RMSE on evaluation set",
             Inches(0.85), Inches(3.75), Inches(5.5), Inches(0.4),
             size=12, color=TEXT_GRAY)
    add_text(s, "R² = 0.00", Inches(0.85), Inches(4.15),
             Inches(5.5), Inches(0.4),
             size=14, color=TEXT_DARK)
    add_text(s, "No variance explained",
             Inches(0.85), Inches(4.6), Inches(5.5), Inches(0.4),
             size=12, italic=True, color=TEXT_GRAY)

    # Arrow
    add_text(s, "▶", Inches(6.55), Inches(3.8),
             Inches(0.5), Inches(0.7),
             size=40, color=HSLU_RED, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(7.05), Inches(1.85), Inches(5.85), Inches(4.7),
             fill=DARK, rounded=True)
    add_text(s, "LightGBM + KNN features",
             Inches(7.3), Inches(2.0), Inches(5.35), Inches(0.4),
             size=13, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s, "CHF 393",
             Inches(7.3), Inches(2.55), Inches(5.35), Inches(1.2),
             size=66, bold=True, color=HSLU_RED, font=FONT_HEAD)
    add_text(s, "RMSE on evaluation set",
             Inches(7.3), Inches(3.75), Inches(5.35), Inches(0.4),
             size=12, color=RGBColor(0xAA, 0xAA, 0xAE))
    add_text(s, "R² = 0.751",
             Inches(7.3), Inches(4.15), Inches(5.35), Inches(0.4),
             size=14, color=TEXT_LIGHT, bold=True)
    add_text(s, "More than 50% improvement",
             Inches(7.3), Inches(4.6), Inches(5.35), Inches(0.4),
             size=12, italic=True, color=GREEN_OK, bold=True)


def slide_actual_predicted(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Actual vs predicted",
                "Best model: LightGBM with KNN-distance features",
                page=page, total=total)
    add_image_or_placeholder(s, "actual_vs_predicted.png",
                              Inches(0.6), Inches(1.85),
                              Inches(7.0), Inches(5.0),
                              caption="Points should cluster around the red diagonal.")
    add_text(s, "What the scatter shows",
             Inches(7.9), Inches(1.85), Inches(5.0), Inches(0.4),
             size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "Strong correlation between actual and predicted price.",
        "Cluster around the diagonal up to ~3'500 CHF.",
        "Tail-end apartments (luxury) show larger spread.",
        "No systematic bias in the central range.",
    ], Inches(7.9), Inches(2.4), Inches(5.2), Inches(3.5), size=14)


def slide_residuals(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Residual analysis",
                "Residual = actual rent minus predicted rent",
                page=page, total=total)
    add_image_or_placeholder(s, "residuals.png",
                              Inches(0.6), Inches(1.85),
                              Inches(7.5), Inches(4.6),
                              caption="Residual scatter + histogram.")
    add_text(s, "Key observations", Inches(8.4), Inches(1.85),
             Inches(4.6), Inches(0.4),
             size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "Funnel shape: errors grow with predicted price.",
        "Distribution is heavier-tailed than a normal.",
        "Shapiro-Wilk test rejects normality.",
        "Confirms heteroscedastic, right-skewed errors.",
    ], Inches(8.4), Inches(2.4), Inches(4.6), Inches(3.6), size=13)


def slide_price_band_errors(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Error per price band",
                "The most actionable finding of the project",
                page=page, total=total)
    bands = [
        ("Cheap",          "CHF 216", 30, ACCENT_BLUE),
        ("Medium-low",     "CHF 176", 25, ACCENT_BLUE),
        ("Medium-high",    "CHF 221", 31, ACCENT_BLUE),
        ("Expensive",      "CHF 445", 62, HSLU_RED),
    ]
    x_left = Inches(0.6)
    bar_top = Inches(2.4)
    bar_h = Inches(0.7)
    gap = Inches(0.2)
    for i, (label, value, pct, color) in enumerate(bands):
        y = bar_top + i * (bar_h + gap)
        width_in = 0.06 + pct * 0.10
        add_rect(s, x_left, y, Inches(2.6), bar_h,
                 fill=CARD_BG, rounded=False)
        add_text(s, label, x_left + Inches(0.2), y + Inches(0.18),
                 Inches(2.4), Inches(0.4),
                 size=14, bold=True, color=TEXT_DARK)
        add_rect(s, x_left + Inches(2.7), y, Inches(width_in), bar_h,
                 fill=color, rounded=False)
        add_text(s, value,
                 x_left + Inches(2.7 + width_in + 0.15),
                 y + Inches(0.15),
                 Inches(2.0), Inches(0.45),
                 size=15, bold=True, color=color, font=FONT_HEAD)
    add_rect(s, Inches(9.5), Inches(2.3), Inches(3.4), Inches(4.4),
             fill=DARK, rounded=True)
    add_text(s, "Why?", Inches(9.7), Inches(2.45),
             Inches(3.0), Inches(0.4),
             size=14, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s,
             "Expensive listings are underrepresented in the training set.",
             Inches(9.7), Inches(2.85), Inches(3.0), Inches(1.0),
             size=13, color=TEXT_LIGHT)
    add_text(s,
             "And many of their price drivers (lake view, penthouse, floor, finish) are not in our features.",
             Inches(9.7), Inches(3.95), Inches(3.0), Inches(2.6),
             size=13, color=TEXT_LIGHT)
    add_text(s,
             "Errors on the top quartile are roughly twice as large as in the other bands.",
             Inches(0.6), Inches(6.5), Inches(12.0), Inches(0.4),
             size=12, italic=True, color=TEXT_GRAY)


def slide_feature_importance(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Feature importance",
                "Four methods agree: location and area carry the signal",
                page=page, total=total)
    add_image_or_placeholder(s, "feature_importance_heatmap.png",
                              Inches(0.6), Inches(1.85),
                              Inches(7.5), Inches(4.8),
                              caption="RF impurity, LightGBM splits, permutation importance and SHAP.")
    add_text(s, "Top drivers", Inches(8.4), Inches(1.85),
             Inches(4.6), Inches(0.4),
             size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_bullets(s, [
        "area: dominant predictor across all methods",
        "knn_price_mean / knn_price_median: location proxy",
        "east, north: encode the location premium",
        "area_per_room beats raw rooms count",
    ], Inches(8.4), Inches(2.4), Inches(4.6), Inches(3.5), size=13)
    add_rect(s, Inches(8.4), Inches(5.5), Inches(4.5), Inches(1.2),
             fill=CARD_BG, rounded=True)
    add_text(s, "Spearman rank correlation",
             Inches(8.55), Inches(5.6), Inches(4.3), Inches(0.4),
             size=12, italic=True, color=TEXT_GRAY)
    add_text(s, "> 0.85",
             Inches(8.55), Inches(5.95), Inches(4.3), Inches(0.6),
             size=26, bold=True, color=GREEN_OK, font=FONT_HEAD)
    add_text(s, "across all four methods",
             Inches(8.55), Inches(6.35), Inches(4.3), Inches(0.3),
             size=10, color=TEXT_GRAY)


# ---- Section 5: Demo --------------------------------------------------------
def slide_demo(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Demo: Streamlit app",
                "End-to-end pipeline behind one address field",
                page=page, total=total)
    # Two-column: left flow, right output
    flow = [
        ("1", "Enter the apartment address"),
        ("2", "Look up EGID and coordinates (GeoAdmin)"),
        ("3", "Fetch GWR details and dwellings list"),
        ("4", "Pick the right apartment from the dropdown"),
        ("5", "Enrich with swisstopo geo-data"),
        ("6", "RentPredictor.predict(...)"),
    ]
    add_text(s, "Live workflow", Inches(0.6), Inches(1.85),
             Inches(6.0), Inches(0.4),
             size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    for i, (num, label) in enumerate(flow):
        y = Inches(2.4 + i * 0.65)
        add_rect(s, Inches(0.6), y, Inches(0.6), Inches(0.55),
                 fill=HSLU_RED, rounded=True)
        add_text(s, num, Inches(0.6), y + Inches(0.07),
                 Inches(0.6), Inches(0.4),
                 size=15, bold=True, color=TEXT_LIGHT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, label, Inches(1.35), y + Inches(0.1),
                 Inches(5.0), Inches(0.45),
                 size=13, color=TEXT_DARK)

    add_rect(s, Inches(7.0), Inches(1.85), Inches(5.85), Inches(4.7),
             fill=DARK, rounded=True)
    add_text(s, "Estimated rent",
             Inches(7.25), Inches(2.0), Inches(5.4), Inches(0.4),
             size=14, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s, "CHF 1'987",
             Inches(7.25), Inches(2.55), Inches(5.4), Inches(1.4),
             size=58, bold=True, color=HSLU_RED, font=FONT_HEAD)
    add_text(s, "for: 80 m², 3 rooms, 8002 Zürich, EGID 12345",
             Inches(7.25), Inches(3.95), Inches(5.4), Inches(0.4),
             size=12, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_rect(s, Inches(7.25), Inches(4.5), Inches(1.5), Inches(0.04),
             fill=HSLU_RED)
    add_text(s, "Includes auto-filled GWR data",
             Inches(7.25), Inches(4.65), Inches(5.4), Inches(0.4),
             size=12, color=TEXT_LIGHT)
    add_text(s, "Editable before prediction",
             Inches(7.25), Inches(5.05), Inches(5.4), Inches(0.4),
             size=12, color=TEXT_LIGHT)
    add_text(s, "Disclaimer: advisory tool, not legally binding",
             Inches(7.25), Inches(5.85), Inches(5.4), Inches(0.4),
             size=10, italic=True, color=RGBColor(0xAA, 0xAA, 0xAE))


# ---- Section 6: Conclusions -------------------------------------------------
def slide_limitations(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Honest limitations", page=page, total=total)
    cards = [
        ("Snapshot dataset",
         "One scrape on 13 April 2026. Market shifts are not reflected."),
        ("Urban coverage bias",
         "Zurich, Geneva, Basel are overrepresented; rural cantons are thin."),
        ("Expensive apartments",
         "Top quartile MAE is roughly twice the rest. Luxury attributes are missing."),
        ("Single platform",
         "All data from rentumo.ch. Adding a second source is the obvious next step."),
        ("Local-only demo",
         "Streamlit runs locally; no hosting in DSPRO1 scope."),
        ("No textual features",
         "Description field is stored but not yet used for feature extraction."),
    ]
    for i, (head, body) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(1.95 + row * 2.45)
        add_rect(s, x, y, Inches(4.0), Inches(2.1),
                 fill=CARD_BG, rounded=True)
        add_rect(s, x, y, Inches(4.0), Inches(0.45),
                 fill=AMBER)
        add_text(s, head, x + Inches(0.2), y + Inches(0.05),
                 Inches(3.7), Inches(0.4),
                 size=14, bold=True, color=TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, body, x + Inches(0.2), y + Inches(0.55),
                 Inches(3.6), Inches(1.5),
                 size=12, color=TEXT_GRAY)


def slide_future_work(prs, total, page):
    s = add_blank_slide(prs)
    page_chrome(s, "Future work, in priority order",
                "Each step targets a specific limitation from the previous slide",
                page=page, total=total)
    items = [
        ("01", "Text features from descriptions",
         "Extract balcony, parking, renovation year, floor, view from the existing description field. Likely the highest-impact lever for the expensive band.",
         HSLU_RED),
        ("02", "Second rental platform",
         "Pull homegate.ch or immoscout24.ch listings. Reduces urban bias and roughly doubles the training set.",
         ACCENT_BLUE),
        ("03", "Prediction intervals",
         "Move from point predictions to CHF X ± Y using conformal prediction (MAPIE). Honest and easy to read.",
         GREEN_OK),
        ("04", "Lightweight hosting",
         "Deploy the Streamlit demo on Streamlit Cloud or Render. Natural first deliverable for DSPRO2.",
         AMBER),
    ]
    for i, (num, head, body, color) in enumerate(items):
        y = Inches(1.85 + i * 1.25)
        add_rect(s, Inches(0.6), y, Inches(12.2), Inches(1.1),
                 fill=CARD_BG, rounded=True)
        add_rect(s, Inches(0.6), y, Inches(1.0), Inches(1.1),
                 fill=color, rounded=False)
        add_text(s, num, Inches(0.6), y,
                 Inches(1.0), Inches(1.1),
                 size=28, bold=True, color=TEXT_LIGHT, font=FONT_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, head, Inches(1.8), y + Inches(0.1),
                 Inches(10.8), Inches(0.4),
                 size=15, bold=True, color=TEXT_DARK, font=FONT_HEAD)
        add_text(s, body, Inches(1.8), y + Inches(0.5),
                 Inches(10.8), Inches(0.6),
                 size=12, color=TEXT_GRAY)


def slide_conclusions(prs, total, page):
    s = add_blank_slide(prs, dark=True)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=HSLU_RED)
    add_text(s, "What we take away",
             Inches(0.7), Inches(0.6), Inches(11.0), Inches(0.6),
             size=32, bold=True, color=TEXT_LIGHT, font=FONT_HEAD)
    add_rect(s, Inches(0.7), Inches(1.35), Inches(1.0), Inches(0.04),
             fill=HSLU_RED)
    facts = [
        ("> 50%", "RMSE reduction vs the mid-term baseline"),
        ("0.751", "R² on the held-out evaluation set"),
        ("CHF 393", "average RMSE of the best model"),
        ("11", "figures generated from the notebook"),
    ]
    for i, (big, small) in enumerate(facts):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.85 + row * 1.7)
        add_text(s, big, x, y,
                 Inches(5.8), Inches(0.9),
                 size=54, bold=True, color=HSLU_RED, font=FONT_HEAD)
        add_text(s, small, x, y + Inches(0.95),
                 Inches(5.8), Inches(0.5),
                 size=14, color=TEXT_LIGHT)
    add_text(s,
             "Tree-based models clearly beat linear regression. The hypothesis is confirmed.",
             Inches(0.7), Inches(5.65), Inches(12.0), Inches(0.5),
             size=16, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))
    add_text(s,
             "The expensive-apartment gap is the most promising next target.",
             Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.5),
             size=16, italic=True, color=RGBColor(0xCC, 0xCC, 0xCE))


def slide_thank_you(prs, total, page):
    s = add_blank_slide(prs, dark=True)
    add_rect(s, Inches(0), Inches(6.7), SLIDE_W, Inches(0.05), fill=HSLU_RED)
    add_text(s, "Thank you for your attention.",
             Inches(0.7), Inches(2.3), Inches(12.0), Inches(0.9),
             size=44, bold=True, color=TEXT_LIGHT, font=FONT_HEAD)
    add_rect(s, Inches(0.7), Inches(3.3), Inches(0.7), Inches(0.04),
             fill=HSLU_RED)
    add_text(s,
             "Questions and feedback welcome.",
             Inches(0.7), Inches(3.5), Inches(12.0), Inches(0.5),
             size=18, italic=True, color=RGBColor(0xBB, 0xBB, 0xBE))
    add_text(s, "Team 8", Inches(0.7), Inches(4.8),
             Inches(6.0), Inches(0.4),
             size=14, bold=True, color=TEXT_LIGHT)
    add_text(s, "Elias Martinelli", Inches(0.7), Inches(5.2),
             Inches(6.0), Inches(0.4), size=13, color=TEXT_LIGHT)
    add_text(s, "Timo Schlumpf", Inches(0.7), Inches(5.55),
             Inches(6.0), Inches(0.4), size=13, color=TEXT_LIGHT)
    add_text(s, "Supervisor", Inches(7.5), Inches(4.8),
             Inches(4.5), Inches(0.4),
             size=14, bold=True, color=TEXT_LIGHT)
    add_text(s, "Elena Nazarenko", Inches(7.5), Inches(5.2),
             Inches(4.5), Inches(0.4), size=13, color=TEXT_LIGHT)
    add_text(s, "Repository", Inches(7.5), Inches(5.7),
             Inches(4.5), Inches(0.4),
             size=12, italic=True, color=RGBColor(0xBB, 0xBB, 0xBE))
    add_text(s, "github.com/wadafacc/dspro1", Inches(7.5), Inches(6.0),
             Inches(4.5), Inches(0.4),
             size=13, color=HSLU_RED)


# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Total = number of slides we will create.  Update if you add/remove slides.
    TOTAL = 22

    slide_title(prs)
    slide_agenda(prs, TOTAL)

    slide_section_divider(prs, 1, "The problem",
                           "Why a fair Swiss rent is hard to judge",
                           TOTAL, 3)
    slide_problem(prs, TOTAL, 4)
    slide_relevance(prs, TOTAL, 5)
    slide_hypothesis(prs, TOTAL, 6)

    slide_section_divider(prs, 2, "Data and pipeline",
                           "From 25'002 URLs to about 4'500 modelling rows",
                           TOTAL, 7)
    slide_data_source(prs, TOTAL, 8)
    slide_scraper(prs, TOTAL, 9)
    slide_enrichment(prs, TOTAL, 10)
    slide_data_funnel(prs, TOTAL, 11)

    slide_section_divider(prs, 3, "Modelling approach",
                           "Six models on the same evaluation protocol",
                           TOTAL, 12)
    slide_methods(prs, TOTAL, 13)
    slide_eval_protocol(prs, TOTAL, 14)

    slide_section_divider(prs, 4, "Results and analysis",
                           "What works, what does not, and why",
                           TOTAL, 15)
    slide_baseline(prs, TOTAL, 16)
    slide_model_comparison(prs, TOTAL, 17)
    slide_hypothesis_confirmed(prs, TOTAL, 18)
    slide_actual_predicted(prs, TOTAL, 19)
    slide_residuals(prs, TOTAL, 20)
    slide_price_band_errors(prs, TOTAL, 21)
    slide_feature_importance(prs, TOTAL, 22)

    slide_section_divider(prs, 5, "Demo",
                           "Streamlit app with address lookup",
                           TOTAL, 23)
    slide_demo(prs, TOTAL, 24)

    slide_section_divider(prs, 6, "Wrap-up",
                           "Limitations, next steps, take-aways",
                           TOTAL, 25)
    slide_limitations(prs, TOTAL, 26)
    slide_future_work(prs, TOTAL, 27)
    slide_conclusions(prs, TOTAL, 28)
    slide_thank_you(prs, TOTAL, 29)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
